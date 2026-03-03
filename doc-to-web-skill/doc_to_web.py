#!/usr/bin/env python3
"""
AI-Powered Document to Webpage Converter
Uses Code Puppy's agent system for intelligent content structuring.
"""

import argparse
import os
import sys
import shutil
import json
import subprocess
from pathlib import Path
from jinja2 import Environment, FileSystemLoader

# Import extractors
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from docx import Document
    import pdfplumber
    from PIL import Image
    import io
except ImportError as e:
    print(f"❌ Error: Missing dependency {e.name}")
    print("Run: uv pip install -r doc-to-web-skill/requirements.txt")
    sys.exit(1)

def invoke_code_puppy_agent(prompt, session_id=None):
    """
    Invokes Code Puppy as a sub-agent to analyze content.
    This uses the Code Puppy CLI to call the agent system.
    """
    cmd = ['code-puppy', '--non-interactive']
    
    if session_id:
        cmd.extend(['--session', session_id])
    
    cmd.append(prompt)
    
    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120
        )
        
        if result.returncode == 0:
            return result.stdout.strip()
        else:
            raise Exception(f"Agent invocation failed: {result.stderr}")
    
    except subprocess.TimeoutExpired:
        raise Exception("Agent invocation timed out")
    except FileNotFoundError:
        raise Exception("code-puppy command not found. Make sure Code Puppy is installed.")

def ai_structure_content(flat_content):
    """
    Uses Code Puppy's AI agent to intelligently structure flat content.
    This is the smart way - let the LLM understand the content!
    """
    print("\n🧠 Using Code Puppy AI to structure content...")
    
    # Prepare the raw content for AI analysis
    content_items = []
    for i, item in enumerate(flat_content[:150]):  # Limit to avoid overwhelming the agent
        if item['type'] == 'text':
            # Truncate very long text
            text = item['value'][:300] + "..." if len(item['value']) > 300 else item['value']
            content_items.append(f"[{i}] TEXT: {text}")
        elif item['type'] == 'image':
            content_items.append(f"[{i}] IMAGE: {item['src']}")
    
    raw_content_str = "\n".join(content_items)
    
    if len(content_items) == 0:
        print("⚠️  Warning: No content to analyze!")
        return [{'title': 'Empty Document', 'blocks': []}]
    
    # AI Analysis Prompt for Code Puppy Agent
    analysis_prompt = f"""You are an expert content strategist analyzing a document for webpage conversion.

RAW DOCUMENT CONTENT:
{raw_content_str}

---

ANALYZE THIS CONTENT AND RETURN STRUCTURED JSON:

Your task:
1. Identify 3-8 logical SECTIONS based on topics/themes
2. For each section, create:
   - title: Short, punchy (max 8 words)
   - lead: Optional 1-2 sentence subtitle  
   - blocks: Array of content blocks

BLOCK TYPES:
- "text": Regular paragraph (use for body content)
- "achievements": Bullet points/accomplishments → bordered cards (use for lists with • or achievements)
- "emphasis": Short impactful statements → centered italic (use for quotes/key messages)
- "two_col": Text + single image → side-by-side layout
- "photo_gallery": 3+ images → grid layout (layouts: "mosaic", "5grid", or "strip")
- "image": Single standalone image
- "divider": Visual section break

RULES:
✅ Bullet points/lists with • → "achievements" block
✅ 3+ consecutive images → "photo_gallery" block  
✅ Text near single image → "two_col" block
✅ Short quotes/impactful statements → "emphasis" block
✅ Group related content logically into sections
✅ For achievements, bold key terms: "<strong>Term:</strong> description"
✅ For emphasis, use <em> tags: "This was <em>transformative</em>"

JSON FORMAT:
{{
  "sections": [
    {{
      "title": "Associate Engagement",
      "lead": "Driving involvement and satisfaction in the workplace",
      "blocks": [
        {{"type": "text", "content": "In FY26, we elevated engagement..."}},
        {{"type": "achievements", "items": ["<strong>Performance:</strong> Enhanced for 887K associates", "<strong>Talent Mobility:</strong> Reached 7,950 associates"]}},
        {{"type": "photo_gallery", "layout": "mosaic", "images": ["[10]", "[11]", "[12]", "[13]", "[14]", "[15]"]}},
        {{"type": "emphasis", "content": "This was a year of <em>transformation</em> and <em>trust</em>."}}
      ]
    }}
  ]
}}

For images, reference by index like "[10]". We'll replace with actual paths.

RETURN ONLY THE JSON - NO MARKDOWN, NO EXPLANATION, JUST VALID JSON.
"""
    
    try:
        # Invoke Code Puppy agent for analysis
        print("   → Sending content to AI agent...")
        ai_response = invoke_code_puppy_agent(analysis_prompt)
        print("   ✓ Received AI analysis")
        
        # Parse JSON (remove markdown code blocks if present)
        if ai_response.startswith("```"):
            ai_response = ai_response.split("```")[1]
            if ai_response.startswith("json"):
                ai_response = ai_response[4:]
            ai_response = ai_response.strip()
        
        structured_data = json.loads(ai_response)
        sections = structured_data['sections']
        
        # Post-process: Replace image index references with actual paths
        image_map = {}
        for i, item in enumerate(flat_content):
            if item['type'] == 'image':
                image_map[f"[{i}]"] = item['src']
        
        for section in sections:
            for block in section['blocks']:
                if block['type'] == 'photo_gallery':
                    block['images'] = [image_map.get(ref, ref) for ref in block['images']]
                elif block['type'] == 'two_col' and 'image_ref' in block:
                    block['image_src'] = image_map.get(block['image_ref'], block['image_ref'])
                    del block['image_ref']
                elif block['type'] == 'image' and 'src_ref' in block:
                    block['src'] = image_map.get(block['src_ref'], block['src_ref'])
                    del block['src_ref']
        
        print(f"   ✓ Created {len(sections)} sections with AI intelligence\n")
        return sections
        
    except Exception as e:
        print(f"\n⚠️  AI structuring failed: {e}")
        print("   → Falling back to simple structure...\n")
        # Fallback: create one section with all content as text blocks
        fallback_blocks = []
        for item in flat_content:
            if item['type'] == 'text':
                fallback_blocks.append({'type': 'text', 'content': item['value']})
            elif item['type'] == 'image':
                fallback_blocks.append({'type': 'image', 'src': item['src']})
        
        return [{
            'title': 'Document Content',
            'lead': 'AI analysis unavailable - displaying raw content',
            'blocks': fallback_blocks
        }]

def extract_from_pptx(filepath, assets_dir):
    """Extract raw text and images from PowerPoint. No structure detection - AI will handle that."""
    print(f"  --> Extracting from PowerPoint: {filepath}")
    prs = Presentation(filepath)
    flat_content = []
    image_count = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        # Sort shapes top-to-bottom, left-to-right
        shapes = sorted(slide.shapes, key=lambda s: (s.top if hasattr(s, 'top') else 0, s.left if hasattr(s, 'left') else 0))
        
        for shape in shapes:
            # Extract text
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    for para in text.split('\n'):
                        if para.strip():
                            flat_content.append({'type': 'text', 'value': para.strip()})
            
            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext = image.ext
                    filename = f"slide_{slide_idx+1}_img_{image_count}.{ext}"
                    image_path = assets_dir / filename
                    
                    with open(image_path, 'wb') as f:
                        f.write(image.blob)
                    
                    rel_path = f"assets/{filename}"
                    flat_content.append({'type': 'image', 'src': rel_path})
                    image_count += 1
                except:
                    pass
                
    return flat_content

def extract_from_docx(filepath, assets_dir):
    """Extract raw text and images from Word doc. No structure detection - AI will handle that."""
    print(f"  --> Extracting from Word Doc: {filepath}")
    doc = Document(filepath)
    flat_content = []
    
    # Extract all text paragraphs in order
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            flat_content.append({'type': 'text', 'value': text})

    # Second pass: Extract images and insert them in sequence
    # Word doesn't make it easy to map images to exact positions,
    # so we extract them and insert after the first few paragraphs
    image_count = 0
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_bytes = rel.target_part.blob
                filename = f"docx_img_{image_count}.png"
                image_path = assets_dir / filename
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                images.append(f"assets/{filename}")
                image_count += 1
            except: 
                pass
    
    # Insert images intelligently: after every ~5 text blocks or at section breaks
    if images:
        enhanced_content = []
        text_block_count = 0
        images_used = 0
        
        for item in flat_content:
            enhanced_content.append(item)
            
            if item['type'] == 'heading' and images_used < len(images):
                # Add 1-2 images after headings
                for _ in range(min(2, len(images) - images_used)):
                    enhanced_content.append({'type': 'image', 'src': images[images_used]})
                    images_used += 1
            elif item['type'] == 'text':
                text_block_count += 1
                # Every 8 text blocks, insert an image
                if text_block_count % 8 == 0 and images_used < len(images):
                    enhanced_content.append({'type': 'image', 'src': images[images_used]})
                    images_used += 1
        
        # Append any remaining images at the end
        while images_used < len(images):
            enhanced_content.append({'type': 'image', 'src': images[images_used]})
            images_used += 1
            
        flat_content = enhanced_content

    return flat_content

def extract_from_pdf(filepath, assets_dir):
    """Extract raw text from PDF. No structure detection - AI will handle that."""
    print(f"  --> Extracting from PDF: {filepath}")
    flat_content = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    for line in text.split('\n'):
                        if line.strip():
                            flat_content.append({'type': 'text', 'value': line.strip()})
    except:
        pass
    return flat_content

def main():
    parser = argparse.ArgumentParser(description="AI-Powered Document to Webpage Converter")
    parser.add_argument("input_file", help="Path to input file (.pptx, .docx, .pdf)")
    parser.add_argument("--output-dir", default="output", help="Directory for output HTML and assets")
    parser.add_argument("--extract-only", action="store_true", 
                        help="Only extract raw content to JSON (for manual AI processing)")
    parser.add_argument("--structure-json", type=str,
                        help="Skip extraction, use provided structured JSON file")
    args = parser.parse_args()
    
    # Create output directory
    output_dir = Path(args.output_dir)
    assets_dir = output_dir / "assets"
    
    # If using pre-structured JSON, skip extraction
    if args.structure_json:
        print(f"\n📥 Loading pre-structured JSON: {args.structure_json}")
        with open(args.structure_json, 'r') as f:
            data = json.load(f)
        sections = data['sections']
        flat_content = data.get('raw_content', [])
        input_path = Path(args.input_file) if Path(args.input_file).exists() else Path(args.structure_json)
        
        # Ensure output directory exists
        if not output_dir.exists():
            os.makedirs(output_dir, exist_ok=True)
        if not assets_dir.exists():
            os.makedirs(assets_dir, exist_ok=True)
        
        # Copy assets from extraction directory if they exist
        source_json_dir = Path(args.structure_json).parent
        source_assets = source_json_dir / "assets"
        if source_assets.exists():
            print(f"   → Copying assets from {source_assets}")
            shutil.copytree(source_assets, assets_dir, dirs_exist_ok=True)
    else:
        # Extract content from document
        input_path = Path(args.input_file)
        if not input_path.exists():
            print(f"❌ Error: File {input_path} not found.")
            sys.exit(1)
        
        # Clean and recreate output directory
        if output_dir.exists():
            shutil.rmtree(output_dir)
        
        os.makedirs(assets_dir, exist_ok=True)
        
        print(f"\n📄 Processing: {input_path.name}")
        print("=" * 60)
        
        ext = input_path.suffix.lower()
        flat_content = []
        
        if ext == '.pptx':
            flat_content = extract_from_pptx(input_path, assets_dir)
        elif ext == '.docx':
            flat_content = extract_from_docx(input_path, assets_dir)
        elif ext == '.pdf':
            flat_content = extract_from_docx(input_path, assets_dir)
        else:
            print(f"❌ Unsupported file type: {ext}")
            sys.exit(1)
        
        print(f"   ✓ Extracted {len([x for x in flat_content if x['type'] == 'text'])} text blocks")
        print(f"   ✓ Extracted {len([x for x in flat_content if x['type'] == 'image'])} images")
        
        # If extract-only mode, save raw content and exit
        if args.extract_only:
            raw_output = output_dir / "raw_content.json"
            with open(raw_output, 'w') as f:
                json.dump({'raw_content': flat_content}, f, indent=2)
            print(f"\n📝 Raw content saved to: {raw_output}")
            print("   → Now use Code Puppy to analyze this content and create structured JSON")
            return
        
        # Use AI to structure content intelligently
        sections = ai_structure_content(flat_content)
    
    # Extract hero image (first image found)
    hero_image = None
    for section in sections:
        for block in section['blocks']:
            if block['type'] == 'image':
                hero_image = block['src']
                break
            if block['type'] == 'two_col':
                hero_image = block['image_src']
                break
            if block['type'] == 'photo_strip':
                hero_image = block['images'][0]
                break
        if hero_image: break

    # Render HTML
    template_dir = Path(__file__).parent / "templates"
    env = Environment(loader=FileSystemLoader(str(template_dir)))
    env.tests['even'] = lambda x: x % 2 == 0 # Register 'even' test
    
    template = env.get_template("template.html")
    
    html_output = template.render(
        title=input_path.stem.replace('_', ' ').title(),
        sections=sections,
        hero_image=hero_image
    )
    
    output_file = output_dir / "index.html"
    with open(output_file, 'w') as f:
        f.write(html_output)
    
    # Count assets
    num_images = len(list(assets_dir.glob('*'))) if assets_dir.exists() else 0
    
    print("\n" + "=" * 60)
    print("✅ WEBSITE GENERATED SUCCESSFULLY!")
    print("=" * 60)
    print(f"📂 Output: {output_file.absolute()}")
    print(f"🖼️  Images: {num_images}")
    print(f"📑 Sections: {len(sections)}")
    print(f"\n🌐 Open in browser: open {output_file.name}")
    print("=" * 60 + "\n")

if __name__ == "__main__":
    main()
